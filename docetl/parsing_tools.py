import importlib
import io
import os
from functools import wraps
from typing import Any, List, Optional


def with_input_output_key(fn):
    """Decorator that wraps a parser function that takes a single
    string parameter and return list of strings and makes it a full
    parser function that takes an item as a dictionary and return a
    list of dictionaries."""

    @wraps(fn)
    def wrapper(item, input_key="text", output_key="text", **kw):
        if input_key not in item:
            raise ValueError(f"Input key {input_key} not found in item: {item}")
        result = fn(item[input_key], **kw)
        if not isinstance(result, list):
            result = [result]
        return [{output_key: res} for res in result]

    return wrapper


def llama_index_simple_directory_reader(
    item: dict[str, Any], input_key: str = "path"
) -> list[dict[str, Any]]:
    from llama_index.core import SimpleDirectoryReader

    documents = SimpleDirectoryReader(item[input_key]).load_data()
    return [{"text": doc.text, "metadata": doc.metadata} for doc in documents]


def llama_index_wikipedia_reader(
    item: dict[str, Any], input_key: str = "pages"
) -> list[dict[str, Any]]:
    from llama_index.readers.wikipedia import WikipediaReader

    loader = WikipediaReader()
    pages = item[input_key]
    if not isinstance(pages, list):
        pages = [pages]
    documents = loader.load_data(pages=pages, auto_suggest=False)
    # The wikipedia reader does not include the page url in the metadata, which is impractical...
    for name, doc in zip(pages, documents):
        doc.metadata["source"] = "https://en.wikipedia.org/wiki/" + name

    return [{"text": doc.text, "metadata": doc.metadata} for doc in documents]


@with_input_output_key
def whisper_speech_to_text(filename: str) -> list[str]:
    """
    Transcribe speech from an audio file to text using Whisper model via litellm.
    If the file is larger than 25 MB, it's split into 10-minute chunks with 30-second overlap.

    Args:
        filename (str): Path to the mp3 or mp4 file.

    Returns:
        list[str]: Transcribed text.
    """

    from litellm import transcription

    file_size = os.path.getsize(filename)
    if file_size > 25 * 1024 * 1024:  # 25 MB in bytes
        from pydub import AudioSegment

        audio = AudioSegment.from_file(filename)
        chunk_length = 10 * 60 * 1000  # 10 minutes in milliseconds
        overlap = 30 * 1000  # 30 seconds in milliseconds

        chunks = []
        for i in range(0, len(audio), chunk_length - overlap):
            chunk = audio[i : i + chunk_length]
            chunks.append(chunk)

        transcriptions = []

        for i, chunk in enumerate(chunks):
            buffer = io.BytesIO()
            buffer.name = f"temp_chunk_{i}_{os.path.basename(filename)}"
            chunk.export(buffer, format="mp3")
            buffer.seek(0)  # Reset buffer position to the beginning

            response = transcription(model="whisper-1", file=buffer)
            transcriptions.append(response.text)

        return transcriptions
    else:
        with open(filename, "rb") as audio_file:
            response = transcription(model="whisper-1", file=audio_file)

        return [response.text]


@with_input_output_key
def xlsx_to_string(
    filename: str,
    orientation: str = "col",
    col_order: str = None,
    doc_per_sheet: bool = False,
    max_rows_per_str: int = 100
) -> List[str]:
    """
    从 Excel 文件中提取文本内容。支持按列或按行聚合数据，并且可以指定表头锚点来定位数据区域。还支持将每个工作表作为独立文档处理。

    主要功能包括：
        1. 定位表头：通过指定 col_order 中的列名，自动识别表头所在行，并以此为基准提取数据。
        2. 数据聚合：根据 orientation 参数选择按列聚合（每列为一个文本块）或按行聚合（每行为一个文本块）。
        3. 分段控制：通过 max_rows_per_str 参数控制每个文本块包含的数据行数，避免单个文本块过大导致处理困难。
        4. 多工作表支持：通过 doc_per_sheet 参数决定是否将每个工作表作为独立文档处理，适用于需要分开分析不同工作表内容的场景。
        5. 支持本地文件和 URL 输入，方便灵活。
        6. 支持合并单元格处理，将父单元格内容合并到子单元格列名中。

    参数:
        - filename (str): Excel 文件的路径或 URL。
        - orientation (str): 数据聚合方式，"col" 表示按列聚合，"row" 表示按行聚合。默认为 "col"。
        - col_order (str, optional): 用于定位表头的列名字符串，列名之间用逗号分隔，函数会自动寻找包含这些列名的行作为表头。默认为 None，即不使用表头定位。
        - doc_per_sheet (bool): 是否将每个工作表作为独立文档处理。默认为 False。
        - max_rows_per_str (int): 每个文本块包含的最大数据行数，默认为 100。

    返回:
        List[str]: 提取的文本内容列表，每个元素对应一个文本块。
    """
    import openpyxl
    import io, requests

    col_order = [col.strip() for col in col_order.split(",")] if col_order else None
    max_rows_per_str = int(max_rows_per_str)
    print("col_order:", col_order)
    # 1. 获取文件流
    file_stream = None
    if filename.startswith(("http://", "https://")):
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        response = requests.get(filename, headers=headers, timeout=30)
        response.raise_for_status()
        # response.encoding = response.apparent_encoding
        file_stream = io.BytesIO(response.content)
    else:
        with open(filename, 'rb') as f:
            file_stream = io.BytesIO(f.read())
    
    # 2. 加载工作簿（注意：不使用 read_only 模式以支持合并单元格读取）
    wb = openpyxl.load_workbook(file_stream, data_only=True, read_only=False)

    def process_sheet(sheet) -> List[str]:
        # 获取合并单元格信息
        merged_ranges = list(sheet.merged_cells.ranges)
        
        # 构建合并单元格映射：(row, col) -> 父单元格值
        merged_cell_map = {}
        for merged_range in merged_ranges:
            min_row, min_col = merged_range.min_row, merged_range.min_col
            parent_value = sheet.cell(min_row, min_col).value
            for row in range(merged_range.min_row, merged_range.max_row + 1):
                for col in range(merged_range.min_col, merged_range.max_col + 1):
                    merged_cell_map[(row, col)] = parent_value
        
        print(f"DEBUG: 工作表名称: {sheet.title}")
        print(f"DEBUG: 合并单元格数量: {len(merged_ranges)}")
        
        # 读取所有行数据
        all_rows = list(sheet.iter_rows(min_row=1, max_row=sheet.max_row, 
                                        min_col=1, max_col=sheet.max_column))
        
        print(f"DEBUG: 总行数: {len(all_rows)}")
        
        if not all_rows:
            return []
        
        # 1. 寻找表头锚点行
        header_row_idx = None
        if col_order:
            for idx, row in enumerate(all_rows):
                row_str_values = [str(c.value).strip() if c.value is not None else "" for c in row]
                print(f"DEBUG: 第 {idx} 行内容: {row_str_values[:5]}...")  # 只打印前5列
                if any(key in row_str_values for key in col_order):
                    header_row_idx = idx
                    print(f"DEBUG: 找到表头行: {idx}")
                    break
        else:
            header_row_idx = 0
        
        if header_row_idx is None:
            print("DEBUG: 未找到表头行")
            return []
        
        # 2. 检测多行表头
        header_end_idx = header_row_idx
        for check_idx in range(header_row_idx + 1, min(header_row_idx + 5, len(all_rows))):
            check_row = all_rows[check_idx]
            # 检查这一行是否包含表头的一部分（通过检查合并单元格）
            is_header_row = False
            for col_idx in range(len(check_row)):
                actual_row = check_idx + 1
                actual_col = col_idx + 1
                # 如果当前单元格在合并区域内，且合并区域的起始行 <= header_row_idx + 1
                if (actual_row, actual_col) in merged_cell_map:
                    # 查找该合并区域的起始行
                    for mr in merged_ranges:
                        if (actual_row >= mr.min_row and actual_row <= mr.max_row and
                            actual_col >= mr.min_col and actual_col <= mr.max_col):
                            if mr.min_row <= header_row_idx + 1:
                                is_header_row = True
                                break
                if is_header_row:
                    break
            
            if is_header_row:
                header_end_idx = check_idx
                print(f"DEBUG: 检测到多行表头，扩展到第 {check_idx} 行")
            else:
                break
        
        # 3. 构建列名（合并多行表头）
        num_cols = len(all_rows[header_row_idx])
        final_headers = []
        
        for col_idx in range(num_cols):
            header_parts = []
            for row_idx in range(header_row_idx, header_end_idx + 1):
                cell = all_rows[row_idx][col_idx]
                actual_row = row_idx + 1
                actual_col = col_idx + 1
                
                # 获取单元格值（考虑合并单元格）
                if (actual_row, actual_col) in merged_cell_map:
                    cell_value = merged_cell_map[(actual_row, actual_col)]
                else:
                    cell_value = cell.value
                
                if cell_value is not None:
                    cell_str = str(cell_value).strip()
                    if cell_str and cell_str not in header_parts:
                        header_parts.append(cell_str)
            
            # 合并多行表头
            if header_parts:
                final_headers.append(" ".join(header_parts))
            else:
                final_headers.append(f"Col_{col_idx + 1}")
        
        print(f"DEBUG: 最终列名: {final_headers}")
        
        # 4. 提取数据行
        data_start_idx = header_end_idx + 1
        data_rows = []
        for row in all_rows[data_start_idx:]:
            row_values = [cell.value for cell in row]
            # 跳过完全空行
            if any(v is not None for v in row_values):
                data_rows.append(row_values)
        
        print(f"DEBUG: 数据行数: {len(data_rows)}")
        
        if not data_rows:
            return []

        # 5. 数据分段逻辑
        chunked_results = []
        
        for i in range(0, len(data_rows), max_rows_per_str):
            chunk = data_rows[i : i + max_rows_per_str]
            current_chunk_str = []

            if orientation == "col":
                # 按列聚合
                for col_idx, h_name in enumerate(final_headers):
                    values = [str(r[col_idx]) for r in chunk if col_idx < len(r) and r[col_idx] is not None]
                    if values:
                        current_chunk_str.append(f"{h_name}:\n" + "\n".join(values))
                        current_chunk_str.append("")
            else:
                # 按行聚合
                for r in chunk:
                    row_parts = [f"{final_headers[k]}: {r[k]}" 
                                 for k in range(len(final_headers)) 
                                 if k < len(r) and r[k] is not None]
                    if row_parts:
                        current_chunk_str.append(" | ".join(row_parts))

            result_str = "\n".join(current_chunk_str).strip()
            if result_str:
                chunked_results.append(result_str)

        return chunked_results

    try:
        final_list = []
        if doc_per_sheet:
            for sh in wb.worksheets:
                final_list.extend(process_sheet(sh))
        else:
            final_list.extend(process_sheet(wb.active))
        return final_list
    finally:
        wb.close()


@with_input_output_key
def txt_to_string(filename: str) -> list[str]:
    """
    Read the content of a text file and return it as a list of strings (only one element).

    Args:
        filename (str): Path to the txt or md file.

    Returns:
        list[str]: Content of the file as a list of strings.
    """
    with open(filename, "r", encoding="utf-8") as file:
        return [file.read()]


@with_input_output_key
def docx_to_string(
    input_path: str, 
    ocr_for_images: bool = False,
    lang: str = "en",
    lines_per_chunk: int = 100,  # 每块最大行数
    overlap_lines: int = 5       # 语义重叠行数
) -> list[str]:
    """
    Extract text from a Word document. Supports both .doc and .docx formats.
    Can optionally use OCR for embedded images. Supports both local files and URLs.

    Args:
        input_path (str): Path to the docx/doc file or URL.
        ocr_for_images (bool): If True, use PaddleOCR to extract text from images.
        lang (str): Language for OCR (default: "en").
        lines_per_chunk (int): Maximum number of lines per chunk.
        overlap_lines (int): Number of overlapping lines between chunks.

    Returns:
        list[str]: Extracted text from the document.
    """
    import sys
    import tempfile
    import requests
    from docx import Document
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    from docx.parts.image import ImagePart
    from docx.text.paragraph import Paragraph
    import io, os
    
    def convert_doc_to_docx_stream(doc_stream: io.BytesIO) -> io.BytesIO:
        """Convert .doc stream to .docx stream"""
        if sys.platform == 'win32':
            try:
                from win32com import client as wc
                # 保存临时doc文件
                with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as tmp_doc:
                    tmp_doc.write(doc_stream.read())
                    tmp_doc_path = tmp_doc.name
                
                try:
                    word = wc.Dispatch("Word.Application")
                    doc = word.Documents.Open(tmp_doc_path)
                    
                    # 保存为docx
                    tmp_docx_path = tmp_doc_path + 'x'
                    doc.SaveAs(tmp_docx_path, 12)
                    doc.Close()
                    word.Quit()
                    
                    # 读取转换后的文件到流
                    with open(tmp_docx_path, 'rb') as f:
                        docx_stream = io.BytesIO(f.read())
                    
                    # 清理临时文件
                    os.remove(tmp_doc_path)
                    os.remove(tmp_docx_path)
                    
                    return docx_stream
                except Exception as e:
                    if os.path.exists(tmp_doc_path):
                        os.remove(tmp_doc_path)
                    raise ValueError(f"doc文件转换失败: {str(e)}")
            except ImportError:
                raise ValueError("Windows系统需要安装pywin32才能转换doc文件")
    
    def is_image_paragraph(p, doc):
        """判断段落是否是图片"""
        try:
            graph = Paragraph(p, doc)
            images = graph._element.xpath('.//pic:pic')
            for image in images:
                for img_id in image.xpath('.//a:blip/@r:embed'):
                    part = doc.part.related_parts[img_id]
                    if isinstance(part, ImagePart):
                        return True, part
        except:
            pass
        return False, None
    
    def extract_text_from_image(image_part, lang):
        """使用PaddleOCR从图片中提取文字"""
        import numpy as np
        from PIL import Image
        from paddleocr import PaddleOCR
        
        if lang == "cn":
            lang = "ch"
        ocr = PaddleOCR(use_angle_cls=True, lang=lang)
        
        # 将图片数据转换为numpy数组
        img_data = image_part.blob
        img = Image.open(io.BytesIO(img_data)).convert('RGB')
        img_array = np.array(img)
        
        # 执行OCR
        ocr_result = ocr.predict(img_array)
        
        texts = []
        if ocr_result and isinstance(ocr_result, list) and len(ocr_result) > 0:
            result_dict = ocr_result[0]  # 获取第一个(也是唯一一个)字典
            if result_dict.get('rec_texts'):
                # 只提取文本内容
                texts += "\n".join(result_dict['rec_texts']) + "\n"
        
        return "\n".join(texts)
    
    lines_per_chunk = int(lines_per_chunk)
    overlap_lines = int(overlap_lines)
    # 1. 处理网络路径 (URL)
    file_stream = None
    # 1. 获取文件流
    if input_path.startswith(("http://", "https://")):
        headers = {"User-Agent": "Mozilla/5.0..."}
        response = requests.get(input_path, headers=headers, timeout=30)
        file_stream = io.BytesIO(response.content)
        file_ext = os.path.splitext(input_path.split('?')[0])[1].lower()
    else:
        with open(input_path, 'rb') as f:
            file_stream = io.BytesIO(f.read())
        file_ext = os.path.splitext(input_path)[1].lower()

    if file_ext == '.doc':
        file_stream = convert_doc_to_docx_stream(file_stream)

    doc = Document(file_stream)

    # --- 核心逻辑：按行数进行语义分块 (方法3) ---
    all_lines = [] # 存储文档提取出的所有原始行

    for element in doc.element.body:
        if isinstance(element, CT_P):
            is_img, img_part = is_image_paragraph(element, doc)
            if is_img and ocr_for_images:
                try:
                    ocr_text = extract_text_from_image(img_part, lang)
                    if ocr_text:
                        all_lines.extend(ocr_text.split('\n'))
                except:
                    all_lines.append("[图片OCR提取失败]")
            else:
                # 即使是空行也保留，维持文档原貌
                all_lines.append(element.text if element.text else "")
        
        elif isinstance(element, CT_Tbl):
            # 表格处理：每行单元格合并为一行文本
            for row in element.tr_lst:
                row_text = "\t".join([p.text for cell in row.tc_lst for p in cell.p_lst if p.text])
                if row_text:
                    all_lines.append(row_text)

    # 3. 执行分块逻辑
    chunks = []
    start_idx = 0
    total_lines = len(all_lines)

    while start_idx < total_lines:
        # 计算当前块的结束位置
        end_idx = start_idx + lines_per_chunk
        chunk_content = all_lines[start_idx:end_idx]
        
        if chunk_content:
            chunks.append("\n".join(chunk_content))
        
        # 步进：移动 lines_per_chunk 减去 重叠行数
        # 比如从 0-100，下一次从 95-195
        start_idx += (lines_per_chunk - overlap_lines)
        
        # 防止死循环或无效步进
        if (lines_per_chunk - overlap_lines) <= 0:
            break

    return chunks if chunks else [""]


@with_input_output_key
def html_to_string(
    input_path: str,
    ocr_for_images: bool = True,
    lang: str = "en",
    lines_per_chunk: int = 100,
    overlap_lines: int = 5
) -> list[str]:
    """
    Extract text from an HTML file or URL. Can optionally use OCR for embedded images.
    
    Args:
        input_path (str): Path to the HTML file or URL.
        ocr_for_images (bool): If True, use PaddleOCR to extract text from images.
        lang (str): Language for OCR (default: "en").
        lines_per_chunk (int): Maximum number of lines per chunk.
        overlap_lines (int): Number of overlapping lines between chunks.
    
    Returns:
        list[str]: Extracted text from the HTML document.
    """
    import requests
    import io
    from bs4 import BeautifulSoup
    import numpy as np
    from PIL import Image
    from paddleocr import PaddleOCR
    
    def extract_text_from_image_url(img_url: str, base_url: str, lang: str, ocr: PaddleOCR) -> str:
        """从图片URL中使用OCR提取文字"""
        try:
            # 处理相对路径
            if not img_url.startswith(("http://", "https://")):
                from urllib.parse import urljoin
                img_url = urljoin(base_url, img_url)
            
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                "Referer": "/".join(img_url.split("/")[:-1]) # 有些服务器会检查来源
            }
            response = requests.get(img_url, headers=headers, timeout=10)
            response.raise_for_status()
            
            img = Image.open(io.BytesIO(response.content)).convert('RGB')
            img_array = np.array(img)
            
            ocr_result = ocr.predict(img_array)
            
            texts = []
            if ocr_result and isinstance(ocr_result, list) and len(ocr_result) > 0:
                result_dict = ocr_result[0]
                if result_dict.get('rec_texts'):
                    texts = result_dict['rec_texts']

            # 合并所有文本
            full_text = "\n".join(texts) if texts else ""
            
            # 检查是否包含中文字符
            if full_text:
                import re
                has_chinese = bool(re.search(r'[\u4e00-\u9fff]', full_text))
                if not has_chinese:
                    return ""  # 没有中文则返回空字符串
            
            return "\n".join(texts) if texts else "[图片OCR提取失败]"
        except Exception as e:
            print(f"图片OCR提取失败，URL: {img_url}，错误: {str(e)}")
            return f"[图片处理失败: {str(e)}]"
    
    lines_per_chunk = int(lines_per_chunk)
    overlap_lines = int(overlap_lines)
    # 1. 获取HTML内容
    if input_path.startswith(("http://", "https://")):
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Referer": "/".join(input_path.split("/")[:-1])
        }
        response = requests.get(input_path, headers=headers, timeout=30)
        response.raise_for_status()
        response.encoding = response.apparent_encoding
        html_content = response.text
        base_url = input_path
    else:
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        base_url = ""
    
    # 2. 解析HTML
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # 移除script和style标签
    for script in soup(["script", "style"]):
        script.decompose()
    
    all_lines = []
    
    ocr = PaddleOCR(use_angle_cls=True, lang=lang)

    # 3. 遍历HTML元素
    for element in soup.descendants:
        # 处理文本节点
        if element.name is None:  # NavigableString
            text = str(element).strip()
            if text:
                all_lines.append(text)
        
        # 处理图片标签
        elif element.name == 'img' and ocr_for_images:
            img_src = element.get('src')
            if img_src:
                ocr_text = extract_text_from_image_url(img_src, base_url, lang, ocr)
                if ocr_text:
                    all_lines.extend(ocr_text.split('\n'))
        
        # 处理表格
        elif element.name == 'tr':
            cells = element.find_all(['td', 'th'])
            if cells:
                row_text = "\t".join([cell.get_text(strip=True) for cell in cells])
                if row_text:
                    all_lines.append(row_text)
    
    # 4. 执行分块逻辑
    chunks = []
    start_idx = 0
    total_lines = len(all_lines)
    
    while start_idx < total_lines:
        end_idx = start_idx + lines_per_chunk
        chunk_content = all_lines[start_idx:end_idx]
        
        if chunk_content:
            chunks.append("\n".join(chunk_content))
        
        start_idx += (lines_per_chunk - overlap_lines)
        
        if (lines_per_chunk - overlap_lines) <= 0:
            break
    
    return chunks if chunks else [""]

@with_input_output_key
def pptx_to_string(filename: str, doc_per_slide: bool = False) -> list[str]:
    """
    Extract text from a PowerPoint presentation.

    Args:
        filename (str): Path to the pptx file.
        doc_per_slide (bool): If True, return each slide as a separate
            document. If False, return the entire presentation as one document.

    Returns:
        list[str]: Extracted text from the presentation. If doc_per_slide
            is True, each string in the list represents a single slide.
            Otherwise, the list contains a single string with all slides'
            content.
    """
    from pptx import Presentation

    prs = Presentation(filename)
    result = []

    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)

        if doc_per_slide:
            result.append("\n".join(slide_content))
        else:
            result.extend(slide_content)

    if not doc_per_slide:
        result = ["\n".join(result)]

    return result


@with_input_output_key
def azure_di_read(
    filename: str,
    use_url: bool = False,
    include_line_numbers: bool = False,
    include_handwritten: bool = False,
    include_font_styles: bool = False,
    include_selection_marks: bool = False,
    doc_per_page: bool = False,
) -> list[str]:
    """
    > Note to developers: We used [this documentation](https://learn.microsoft.com/en-us/azure/ai-services/document-intelligence/how-to-guides/use-sdk-rest-api?view=doc-intel-4.0.0&tabs=windows&pivots=programming-language-python) as a reference.

    This function uses Azure Document Intelligence to extract text from documents.
    To use this function, you need to set up an Azure Document Intelligence resource:

    1. [Create an Azure account](https://azure.microsoft.com/) if you don't have one
    2. Set up a Document Intelligence resource in the [Azure portal](https://portal.azure.com/#create/Microsoft.CognitiveServicesFormRecognizer)
    3. Once created, find the resource's endpoint and key in the Azure portal
    4. Set these as environment variables:
       - DOCUMENTINTELLIGENCE_API_KEY: Your Azure Document Intelligence API key
       - DOCUMENTINTELLIGENCE_ENDPOINT: Your Azure Document Intelligence endpoint URL

    The function will use these credentials to authenticate with the Azure service.
    If the environment variables are not set, the function will raise a ValueError.

    The Azure Document Intelligence client is then initialized with these credentials.
    It sends the document (either as a file or URL) to Azure for analysis.
    The service processes the document and returns structured information about its content.

    This function then extracts the text content from the returned data,
    applying any specified formatting options (like including line numbers or font styles).
    The extracted text is returned as a list of strings, with each string
    representing either a page (if doc_per_page is True) or the entire document.

    Args:
        filename (str): Path to the file to be analyzed or URL of the document if use_url is True.
        use_url (bool, optional): If True, treat filename as a URL. Defaults to False.
        include_line_numbers (bool, optional): If True, include line numbers in the output. Defaults to False.
        include_handwritten (bool, optional): If True, include handwritten text in the output. Defaults to False.
        include_font_styles (bool, optional): If True, include font style information in the output. Defaults to False.
        include_selection_marks (bool, optional): If True, include selection marks in the output. Defaults to False.
        doc_per_page (bool, optional): If True, return each page as a separate document. Defaults to False.

    Returns:
        list[str]: Extracted text from the document. If doc_per_page is True, each string in the list represents
                   a single page. Otherwise, the list contains a single string with all pages' content.

    Raises:
        ValueError: If DOCUMENTINTELLIGENCE_API_KEY or DOCUMENTINTELLIGENCE_ENDPOINT environment variables are not set.
    """

    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.ai.documentintelligence.models import AnalyzeDocumentRequest
    from azure.core.credentials import AzureKeyCredential

    key = os.getenv("DOCUMENTINTELLIGENCE_API_KEY")
    endpoint = os.getenv("DOCUMENTINTELLIGENCE_ENDPOINT")

    if key is None:
        raise ValueError("DOCUMENTINTELLIGENCE_API_KEY environment variable is not set")
    if endpoint is None:
        raise ValueError(
            "DOCUMENTINTELLIGENCE_ENDPOINT environment variable is not set"
        )

    document_analysis_client = DocumentIntelligenceClient(
        endpoint=endpoint, credential=AzureKeyCredential(key)
    )

    if use_url:
        poller = document_analysis_client.begin_analyze_document(
            "prebuilt-read", AnalyzeDocumentRequest(url_source=filename)
        )
    else:
        with open(filename, "rb") as f:
            poller = document_analysis_client.begin_analyze_document("prebuilt-read", f)

    result = poller.result()

    style_content = []
    content = []

    if result.styles:
        for style in result.styles:
            if style.is_handwritten and include_handwritten:
                handwritten_text = ",".join(
                    [
                        result.content[span.offset : span.offset + span.length]
                        for span in style.spans
                    ]
                )
                style_content.append(f"Handwritten content: {handwritten_text}")

            if style.font_style and include_font_styles:
                styled_text = ",".join(
                    [
                        result.content[span.offset : span.offset + span.length]
                        for span in style.spans
                    ]
                )
                style_content.append(f"'{style.font_style}' font style: {styled_text}")

    for page in result.pages:
        page_content = []

        if page.lines:
            for line_idx, line in enumerate(page.lines):
                if include_line_numbers:
                    page_content.append(f" Line #{line_idx}: {line.content}")
                else:
                    page_content.append(f"{line.content}")

        if page.selection_marks and include_selection_marks:
            # TODO: figure this out
            for selection_mark_idx, selection_mark in enumerate(page.selection_marks):
                page_content.append(
                    f"Selection mark #{selection_mark_idx}: State is '{selection_mark.state}' within bounding polygon "
                    f"'{selection_mark.polygon}' and has a confidence of {selection_mark.confidence}"
                )

        content.append("\n".join(page_content))

    if doc_per_page:
        return style_content + content
    else:
        return [
            "\n\n".join(
                [
                    "\n".join(style_content),
                    "\n\n".join(
                        f"Page {i+1}:\n{page_content}"
                        for i, page_content in enumerate(content)
                    ),
                ]
            )
        ]


@with_input_output_key
def paddleocr_pdf_to_string(
    input_path: str,
    doc_per_page: bool = False,
    ocr_enabled: bool = True,
    lang: str = "ch",
) -> list[str]:
    """
    Extract text and image information from a PDF file using PaddleOCR for image-based PDFs.

    **Note: this is very slow!!**

    Args:
        input_path (str): Path to the input PDF file.
        doc_per_page (bool): If True, return a list of strings, one per page.
            If False, return a single string.
        ocr_enabled (bool): Whether to enable OCR for image-based PDFs.
        lang (str): Language of the PDF file.

    Returns:
        list[str]: Extracted content as a list of formatted strings.
    """
    import fitz
    import numpy as np
    import requests  # 新增：用于处理网络请求
    import io        # 新增：用于处理内存流
    from paddleocr import PaddleOCR
    import logging

    # 1. 处理网络路径 (URL)
    if input_path.startswith(("http://", "https://")):
        print(f"检测到 URL，正在下载文件: {input_path}")
        
        # 添加伪装头，防止 403 错误
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Referer": "/".join(input_path.split("/")[:-1]) # 有些服务器会检查来源
        }
        
        try:
            response = requests.get(input_path, headers=headers, timeout=30)
            response.raise_for_status() 
            # response.encoding = response.apparent_encoding
        except requests.exceptions.HTTPError as e:
            # 如果还是 403，打印更详细的信息
            print(f"下载失败，状态码: {response.status_code}。服务器拒绝了请求。")
            raise e
            
        file_data = io.BytesIO(response.content)
        pdf = fitz.open(stream=file_data, filetype="pdf")
    else:
        # 如果是本地路径，按原逻辑打开
        pdf = fitz.open(input_path)

    # 2. 初始化 OCR (lang="cn" 会报错，这里确保是 "ch")
    if lang == "cn": lang = "ch"
    ocr = PaddleOCR(use_angle_cls=True, lang=lang)

    pdf_content = []

    with pdf:
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            text = page.get_text()
            images = []

            # 提取图像信息 (bbox)
            for img_index, img in enumerate(page.get_images(full=True)):
                try:
                    rect = page.get_image_bbox(img)
                    images.append(f"Image {img_index + 1}: bbox {rect}")
                except Exception:
                    continue

            page_content = f"Page {page_num + 1}:\n"
            page_content += f"Text:\n{text}\n"
            page_content += "Images:\n" + "\n".join(images) + "\n"

            # 3. 如果页面没有文字且开启了 OCR，则转为图片识别
            if not text.strip() and ocr_enabled:
                mat = fitz.Matrix(2, 2) # 提高分辨率提高识别率
                pix = page.get_pixmap(matrix=mat)
                
                # 将渲染的图片转为 numpy 数组供 Paddle 使用
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                    pix.height, pix.width, 3
                )

                ocr_result = ocr.predict(img)
                page_content += "OCR Results:\n"
                # ocr_result 是一个列表,包含一个字典
                if ocr_result and isinstance(ocr_result, list) and len(ocr_result) > 0:
                    result_dict = ocr_result[0]  # 获取第一个(也是唯一一个)字典
                    if result_dict.get('rec_texts'):
                        # 只提取文本内容
                        page_content += "\n".join(result_dict['rec_texts']) + "\n"

            pdf_content.append(page_content)

    if not doc_per_page:
        return ["\n\n".join(pdf_content)]

    return pdf_content


@with_input_output_key
def gptpdf_to_string(
    input_path: str,
    gpt_model: str,
    api_key: str,
    base_url: str,
    verbose: bool = False,
    custom_prompt: dict[str, str] | None = None,
) -> str:
    """
    Parse PDF using GPT to convert the content of a PDF to a markdown format and write it to an output file.

    **Note: pip install gptpdf required**

    Args:
        input_path (str): Path to the input PDF file.
        gpt_model (str): GPT model to be used for parsing.
        api_key (str): API key for GPT service.
        base_url (str): Base URL for the GPT service.
        verbose (bool): If True, will print additional information during parsing.
        custom_prompt (dict[str, str] | None): Custom prompt for the GPT model. See https://github.com/CosmosShadow/gptpdf for more information.

    Returns:
        list[str]: Extracted content as a list of strings.
    """
    import tempfile

    from gptpdf import parse_pdf

    with tempfile.TemporaryDirectory() as temp_dir:
        kwargs = {
            "pdf_path": input_path,
            "output_dir": temp_dir,
            "api_key": api_key,
            "base_url": base_url,
            "model": gpt_model,
            "verbose": verbose,
        }
        if custom_prompt:
            kwargs["prompt"] = custom_prompt

        parsed_content, _ = parse_pdf(
            **kwargs
        )  # The second element is a list of image paths, which we don't need.

        return [parsed_content]


# Define a dictionary mapping function names to their corresponding functions


def get_parser(name: str):
    try:
        entrypoint = importlib.metadata.entry_points(group="docetl.parser")[name]
    except KeyError:
        raise KeyError(f"Unrecognized parser {name}")
    return entrypoint.load()


def get_parsing_tools():
    return [ep.name for ep in importlib.metadata.entry_points(group="docetl.parser")]
